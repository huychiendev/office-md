"""
Subprocess worker for file conversion.
Runs in an isolated process so OOM kills only this process, not the main server.
Usage: python worker.py <job_id> <job_dir> <file_path> <filename>
"""
import sys
import os
import re
import json
import urllib.parse


def _escape_cell(text):
    return text.replace('|', '\\|').replace('\n', ' ')


def _extract_vba_macros(file_path):
    try:
        from oletools.olevba import VBA_Parser
        vba_parser = VBA_Parser(file_path)
        if not vba_parser.detect_vba_macros():
            vba_parser.close()
            return ""
        parts = ["\n## VBA Macros\n"]
        for (filename, stream_path, vba_filename, vba_code_chunk) in vba_parser.extract_macros():
            parts.append(f"### Module: {vba_filename}\n```vba\n{vba_code_chunk}\n```\n")
        vba_parser.close()
        return "\n".join(parts)
    except Exception as e:
        print(f"Error extracting VBA: {e}")
        return ""


def convert_excel_with_formulas(job_id, job_dir, file_path, exclude_hidden_sheets=False):
    import colorsys
    from openpyxl import load_workbook
    from openpyxl.utils import get_column_letter
    from openpyxl.xml.functions import QName, fromstring
    
    def get_theme_colors(wb):
        if not getattr(wb, "loaded_theme", None):
            return []
        try:
            xlmns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
            root = fromstring(wb.loaded_theme)
            theme_el = root.find(QName(xlmns, 'themeElements').text)
            color_scheme = theme_el.find(QName(xlmns, 'clrScheme').text)
            
            theme_names = [
                'lt1', 'dk1', 'lt2', 'dk2', 'accent1', 'accent2', 
                'accent3', 'accent4', 'accent5', 'accent6'
            ]
            
            colors = []
            for name in theme_names:
                element = color_scheme.find(QName(xlmns, name).text)
                if element is None:
                    colors.append("FFFFFF")
                    continue
                sys_clr = element.find(QName(xlmns, 'sysClr').text)
                srgb_clr = element.find(QName(xlmns, 'srgbClr').text)
                
                if srgb_clr is not None:
                    colors.append(srgb_clr.get('val'))
                elif sys_clr is not None:
                    colors.append(sys_clr.get('lastClr') or "FFFFFF")
                else:
                    colors.append("FFFFFF")
            return colors
        except Exception:
            return []

    def apply_tint(rgb_hex, tint):
        if not tint:
            return rgb_hex
        try:
            r = int(rgb_hex[0:2], 16) / 255.0
            g = int(rgb_hex[2:4], 16) / 255.0
            b = int(rgb_hex[4:6], 16) / 255.0
            h, l, s = colorsys.rgb_to_hls(r, g, b)
            if tint < 0:
                l = l * (1.0 + tint)
            else:
                l = l * (1.0 - tint) + tint
            l = max(0.0, min(1.0, l))
            r, g, b = colorsys.hls_to_rgb(h, l, s)
            return f"{int(r * 255):02X}{int(g * 255):02X}{int(b * 255):02X}"
        except Exception:
            return rgb_hex

    def get_cell_color(cell, theme_colors):
        if not cell.fill or cell.fill.fill_type != "solid":
            return None
        color = cell.fill.start_color
        if color.type == "rgb" and isinstance(color.rgb, str):
            rgb = color.rgb[-6:]
            if rgb.lower() != "000000":
                return f"#{rgb}"
        elif color.type == "theme" and color.theme is not None:
            if theme_colors and color.theme < len(theme_colors):
                base_color = theme_colors[color.theme]
                rgb = apply_tint(base_color, color.tint)
                return f"#{rgb}"
        return None

    def get_text_color_for_background(bg_hex):
        if bg_hex.startswith("#"):
            bg_hex = bg_hex[1:]
        try:
            r = int(bg_hex[0:2], 16)
            g = int(bg_hex[2:4], 16)
            b = int(bg_hex[4:6], 16)
            luminance = (0.299 * r + 0.587 * g + 0.114 * b) / 255
            return "#000000" if luminance > 0.5 else "#FFFFFF"
        except Exception:
            return "#FFFFFF"

    style_block = """<style>
  .excel-table { border-collapse: collapse; width: 100%; border: 1px solid rgba(255,255,255,0.1); min-width: 600px; margin-bottom: 20px; }
  .excel-table th { border: 1px solid rgba(255,255,255,0.1); padding: 8px; text-align: center; font-weight: bold; background-color: rgba(255,255,255,0.05); }
  .excel-table td { border: 1px solid rgba(255,255,255,0.1); padding: 8px; vertical-align: top; white-space: pre-wrap; }
  .excel-table .row-idx { font-weight: bold; background-color: rgba(255,255,255,0.05); text-align: center; width: 40px; }
</style>
"""

    wb_data = load_workbook(file_path, data_only=True, read_only=False)
    wb_formula = load_workbook(file_path, data_only=False, read_only=False)
    theme_colors = get_theme_colors(wb_formula)
    parts = [style_block]

    for name in wb_formula.sheetnames:
        ws_d, ws_f = wb_data[name], wb_formula[name]
        
        if exclude_hidden_sheets and getattr(ws_f, 'sheet_state', 'visible') != 'visible':
            continue
        
        # Build merged cells map
        merged_ranges = ws_d.merged_cells.ranges
        merged_map = {}
        skip_cells = set()
        
        for rng in merged_ranges:
            min_col, min_row, max_col, max_row = rng.bounds
            rowspan = max_row - min_row + 1
            colspan = max_col - min_col + 1
            
            for r in range(min_row, max_row + 1):
                for c in range(min_col, max_col + 1):
                    if r == min_row and c == min_col:
                        merged_map[(r, c)] = (rowspan, colspan)
                    else:
                        skip_cells.add((r, c))

        images_by_cell = {}
        img_counter = 0
        max_img_row = -1
        max_img_col = -1
        
        for image in getattr(ws_f, '_images', []):
            try:
                if not hasattr(image, 'anchor') or not hasattr(image.anchor, '_from'):
                    continue
                
                row_idx = image.anchor._from.row
                col_idx = image.anchor._from.col
                
                if row_idx > max_img_row: max_img_row = row_idx
                if col_idx > max_img_col: max_img_col = col_idx
                
                img_counter += 1
                ext = getattr(image, 'format', 'png') or 'png'
                save_name = f"{name}_img_{img_counter}.{ext}"
                save_name = re.sub(r"[^\w\.-]", "_", save_name)
                img_path = os.path.join(job_dir, save_name)
                
                if callable(getattr(image, '_data', None)):
                    blob = image._data()
                else:
                    blob = image._data
                
                with open(img_path, "wb") as f:
                    f.write(blob)
                    
                encoded_filename = urllib.parse.quote(save_name)
                html_img = f"<img src='{encoded_filename}' alt='{save_name}' style='max-height: 100px; display: inline-block; vertical-align: middle; margin: 2px;' />"
                
                if (row_idx, col_idx) not in images_by_cell:
                    images_by_cell[(row_idx, col_idx)] = []
                images_by_cell[(row_idx, col_idx)].append(html_img)
            except Exception as e:
                print(f"Error extracting image from Excel: {e}")
                
        # Determine actual data boundaries based on cell values, formulas and images
        actual_max_row = 1
        actual_max_col = 1
        
        for (r, c), cell in ws_d._cells.items():
            if cell.value is not None:
                if r > actual_max_row:
                    actual_max_row = r
                if c > actual_max_col:
                    actual_max_col = c
                    
        for (r, c), cell in ws_f._cells.items():
            if cell.value is not None:
                if r > actual_max_row:
                    actual_max_row = r
                if c > actual_max_col:
                    actual_max_col = c
                    
        for (img_r, img_c) in images_by_cell.keys():
            r, c = img_r + 1, img_c + 1
            if r > actual_max_row:
                actual_max_row = r
            if c > actual_max_col:
                actual_max_col = c
                
        for rng in ws_d.merged_cells.ranges:
            min_col, min_row, max_col, max_row = rng.bounds
            # Check if this merge range contains any values
            has_val = False
            for r in range(min_row, max_row + 1):
                for c in range(min_col, max_col + 1):
                    if ws_d.cell(row=r, column=c).value is not None:
                        has_val = True
                        break
                if has_val:
                    break
            
            if has_val:
                if max_row > actual_max_row:
                    actual_max_row = max_row
                if max_col > actual_max_col:
                    actual_max_col = max_col

        sheet_title = name
        if getattr(ws_f, 'sheet_state', 'visible') != 'visible':
            sheet_title += " (sheet ẩn)"
        parts.append(f"## {sheet_title}\n")
        
        # Check if there is any data or merge cells or images in the sheet to prevent rendering empty sheets
        has_any_data = False
        for r in range(1, actual_max_row + 1):
            for c in range(1, actual_max_col + 1):
                if ws_d.cell(row=r, column=c).value is not None or (r - 1, c - 1) in images_by_cell or (r, c) in merged_map:
                    has_any_data = True
                    break
            if has_any_data:
                break
                
        if not has_any_data:
            continue
            
        sheet_html = []
        sheet_html.append("<div style='overflow-x: auto; margin-bottom: 20px;'>")
        sheet_html.append("  <table class='excel-table'>")
        
        # Header (Column Letters A, B, C...)
        sheet_html.append("    <thead>")
        sheet_html.append("      <tr>")
        sheet_html.append("        <th class='row-idx'></th>")
        for c in range(1, actual_max_col + 1):
            sheet_html.append(f"        <th>{get_column_letter(c)}</th>")
        sheet_html.append("      </tr>")
        sheet_html.append("    </thead>")
        
        # Body (Rows 1, 2, 3...)
        sheet_html.append("    <tbody>")
        
        for r in range(1, actual_max_row + 1):
            row_has_data = False
            for c in range(1, actual_max_col + 1):
                if ws_d.cell(row=r, column=c).value is not None or (r - 1, c - 1) in images_by_cell or (r, c) in merged_map or (r, c) in skip_cells:
                    row_has_data = True
                    break
            
            if not row_has_data:
                continue
                
            sheet_html.append("      <tr>")
            # Row index cell
            sheet_html.append(f"        <td class='row-idx'>{r}</td>")
            
            for c in range(1, actual_max_col + 1):
                if (r, c) in skip_cells:
                    continue
                    
                cell_d = ws_d.cell(row=r, column=c)
                cell_f = ws_f.cell(row=r, column=c)
                
                val = cell_d.value
                raw = cell_f.value
                
                if isinstance(raw, str) and raw.startswith('='):
                    text = f"{val if val is not None else ''} (`{raw}`)"
                else:
                    text = str(val) if val is not None else ''
                    
                imgs = images_by_cell.get((r - 1, c - 1), [])
                if imgs:
                    text = text + " " + " ".join(imgs) if text else " ".join(imgs)
                
                # Style and merges
                td_style = ""
                color = get_cell_color(cell_d, theme_colors)
                if color:
                    text_color = get_text_color_for_background(color)
                    td_style = f"background-color: {color}; color: {text_color};"
                    
                style_attr = f" style='{td_style}'" if td_style else ""
                
                merge_attrs = ""
                if (r, c) in merged_map:
                    rowspan, colspan = merged_map[(r, c)]
                    if rowspan > 1:
                        merge_attrs += f" rowspan='{rowspan}'"
                    if colspan > 1:
                        merge_attrs += f" colspan='{colspan}'"
                
                sheet_html.append(f"        <td{style_attr}{merge_attrs}>{text}</td>")
                
            sheet_html.append("      </tr>")
            
        sheet_html.append("    </tbody>")
        sheet_html.append("  </table>")
        sheet_html.append("</div>\n")
        
        parts.append("\n".join(sheet_html))

    wb_data.close()
    wb_formula.close()
    return '\n'.join(parts)


def convert_file(job_id, job_dir, file_path, filename, exclude_hidden_sheets=False):
    """Core conversion logic -- runs in isolated process."""
    markdown_text = ""

    if filename.lower().endswith(('.xlsx', '.xlsm', '.xltx', '.xltm')):
        markdown_text = convert_excel_with_formulas(job_id, job_dir, file_path, exclude_hidden_sheets)
    else:
        # Only import MarkItDown here -- keeps main server process lean
        from markitdown import MarkItDown
        md = MarkItDown()
        result = md.convert(file_path)
        markdown_text = result.text_content
        markdown_text = re.sub(r'\bNaN\b', '', markdown_text)
        markdown_text = re.sub(r'Unnamed:\s*\d+', '', markdown_text)

    # VBA macros
    office_exts = ('doc', 'docx', 'docm', 'dot', 'dotx', 'dotm',
                   'xls', 'xlsx', 'xlsm', 'xlsb', 'xlt', 'xltx', 'xltm',
                   'ppt', 'pptx', 'pptm', 'pot', 'potx', 'potm', 'pps', 'ppsx', 'ppsm')
    if filename.lower().endswith(office_exts):
        vba_text = _extract_vba_macros(file_path)
        if vba_text:
            markdown_text += vba_text

    # PPTX images
    if filename.lower().endswith('.pptx'):
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            img_counter = {}
            for slide in prs.slides:
                sorted_shapes = sorted(
                    slide.shapes,
                    key=lambda s: (
                        float("-inf") if not s.top else s.top,
                        float("-inf") if not s.left else s.left,
                    ),
                )
                for shape in sorted_shapes:
                    if not (shape.shape_type == 13 or (shape.shape_type == 14 and hasattr(shape, "image"))):
                        continue
                    try:
                        blob = shape.image.blob
                        ext = shape.image.ext
                    except Exception:
                        continue

                    placeholder_name = re.sub(r"\W", "", shape.name)
                    placeholder = f"{placeholder_name}.jpg"
                    count = img_counter.get(placeholder_name, 0)
                    img_counter[placeholder_name] = count + 1
                    save_name = f"{placeholder_name}_{count}.{ext}" if count > 0 else f"{placeholder_name}.{ext}"

                    img_path = os.path.join(job_dir, save_name)
                    with open(img_path, "wb") as img_f:
                        img_f.write(blob)

                    encoded_filename = urllib.parse.quote(save_name)
                    markdown_text = markdown_text.replace(f"]({placeholder})", f"]({encoded_filename})", 1)
        except Exception as e:
            print(f"Error extracting PPTX images: {e}")

    return markdown_text


def main():
    if len(sys.argv) < 5:
        print("Usage: python worker.py <job_id> <job_dir> <file_path> <filename> [--batch]")
        sys.exit(1)

    job_id, job_dir, file_path, filename = sys.argv[1], sys.argv[2], sys.argv[3], sys.argv[4]
    flags = sys.argv[5:]
    is_batch = "--batch" in flags
    exclude_hidden_sheets = "--exclude-hidden-sheets" in flags

    try:
        markdown_text = convert_file(job_id, job_dir, file_path, filename, exclude_hidden_sheets)

        md_filename = f"{os.path.splitext(filename)[0]}.md"
        md_filepath = os.path.join(job_dir, md_filename)
        os.makedirs(os.path.dirname(md_filepath), exist_ok=True)
        
        with open(md_filepath, "w", encoding="utf-8") as f:
            f.write(markdown_text)

        if not is_batch:
            import shutil
            has_images = any(f.endswith(('.png', '.jpg', '.jpeg', '.gif')) for f in os.listdir(job_dir))
            if has_images:
                zip_path = os.path.join(os.path.dirname(job_dir), f"{job_id}_archive")
                shutil.make_archive(zip_path, 'zip', job_dir)
            with open(os.path.join(job_dir, "success.txt"), "w") as f:
                f.write(md_filename)

    except Exception as e:
        if not is_batch:
            with open(os.path.join(job_dir, "error.txt"), "w", encoding="utf-8") as f:
                f.write(str(e))
        else:
            print(f"Error processing {filename}: {e}", file=sys.stderr)
            sys.exit(1)
    finally:
        if os.path.exists(file_path):
            os.remove(file_path)


if __name__ == "__main__":
    main()
